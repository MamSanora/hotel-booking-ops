import sys
import os
from docx import Document
from pptx import Presentation

def main():
    out_path = r"d:\xampp\htdocs\Project_Sarana\Hotel_Booking_Ops\scratch\extract_output_utf8.txt"
    with open(out_path, "w", encoding="utf-8") as f:
        docx_path = r"d:\xampp\htdocs\Project_Sarana\Hotel_Booking_Ops\z_documentation\defense material\SARANA Complete-Team14 Final_Finalize_v1.docx"
        f.write(f"--- DOCX HEADINGS: {os.path.basename(docx_path)} ---\n")
        try:
            doc = Document(docx_path)
            for p in doc.paragraphs:
                if p.style.name.startswith('Heading'):
                    f.write(f"{p.style.name}: {p.text}\n")
        except Exception as e:
            f.write(f"Error reading docx: {e}\n")
        f.write("\n")

        pptx_path = r"d:\xampp\htdocs\Project_Sarana\Hotel_Booking_Ops\z_documentation\defense material\V2_Main_Dara_Meas_Hotel_Capstone_Defense.pptx"
        f.write(f"--- PPTX CONTENT: {os.path.basename(pptx_path)} ---\n")
        try:
            prs = Presentation(pptx_path)
            for i, slide in enumerate(prs.slides):
                f.write(f"Slide {i+1}:\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.replace('\n', ' ').strip()
                        if text:
                            f.write(f"  - {text}\n")
        except Exception as e:
            f.write(f"Error reading pptx: {e}\n")
        f.write("\n")

if __name__ == '__main__':
    main()
