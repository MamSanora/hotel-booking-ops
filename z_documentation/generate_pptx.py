import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set to 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG_DARK = RGBColor(15, 23, 42)      # #0f172a
    COLOR_BG_CARD = RGBColor(30, 41, 59)      # #1e293b
    COLOR_CARD_BORDER = RGBColor(51, 65, 85)  # #334155
    COLOR_GOLD = RGBColor(200, 169, 110)      # #c8a96e
    COLOR_GOLD_LIGHT = RGBColor(245, 158, 11)  # #f59e0b
    COLOR_WHITE = RGBColor(248, 250, 252)     # #f8fafc
    COLOR_MUTED = RGBColor(148, 163, 184)     # #94a3b8
    COLOR_CYAN = RGBColor(14, 165, 233)       # #0ea5e9
    COLOR_GREEN = RGBColor(16, 185, 129)      # #10b981

    def add_blank_slide(bg_color=COLOR_BG_DARK):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()
        return slide

    def add_header(slide, title_text, category_text="CAPSTONE DEFENSE — NORTON UNIVERSITY"):
        # Category / Kicker
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD
        p.font.name = "Arial"

        # Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.font.name = "Arial"

        # Decorative Gold Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.55), Inches(11.733), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_GOLD
        line.line.fill.background()

    def add_card(slide, left, top, width, height, title="", bg_color=COLOR_BG_CARD, border_color=COLOR_CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()

        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLOR_GOLD
            p.font.name = "Arial"
        return shape

    # ── SLIDE 1: Title Slide ──────────────────────────────────────────────────
    s1 = add_blank_slide()

    # Outer decorative glow frame
    frame = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.6), Inches(12.133), Inches(6.3))
    frame.fill.solid()
    frame.fill.fore_color.rgb = COLOR_BG_CARD
    frame.line.color.rgb = COLOR_GOLD
    frame.line.width = Pt(2)

    # University Badge Header Text
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(0.9), Inches(11.333), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "NORTON UNIVERSITY • FACULTY OF COMPUTER SCIENCE"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.name = "Arial"

    p_sub = tf.add_paragraph()
    p_sub.text = "DEPARTMENT OF COMPUTER SCIENCE (BATCH 27, 2026)"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.size = Pt(11)
    p_sub.font.color.rgb = COLOR_MUTED
    p_sub.font.name = "Arial"

    # Main Project Title Box
    tb_title = s1.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.333), Inches(1.8))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_khmer = tf_title.paragraphs[0]
    p_khmer.text = "ប្រព័ន្ធដំណើរការកក់សណ្ឋាគារ តារាមាស"
    p_khmer.alignment = PP_ALIGN.CENTER
    p_khmer.font.size = Pt(28)
    p_khmer.font.bold = True
    p_khmer.font.color.rgb = COLOR_WHITE
    p_khmer.font.name = "Arial"

    p_eng = tf_title.add_paragraph()
    p_eng.text = "Dara Meas Hotel Booking Operation System"
    p_eng.alignment = PP_ALIGN.CENTER
    p_eng.font.size = Pt(24)
    p_eng.font.bold = True
    p_eng.font.color.rgb = COLOR_GOLD_LIGHT
    p_eng.font.name = "Arial"

    p_tag = tf_title.add_paragraph()
    p_tag.text = "BACHELOR CAPSTONE PROJECT DEFENSE"
    p_tag.alignment = PP_ALIGN.CENTER
    p_tag.font.size = Pt(13)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_CYAN
    p_tag.font.name = "Arial"

    # Divider line
    div = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.95), Inches(6.333), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_GOLD
    div.line.fill.background()

    # Left Team Box
    tb_team = s1.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(5.2), Inches(2.2))
    tf_team = tb_team.text_frame
    tf_team.word_wrap = True
    pt = tf_team.paragraphs[0]
    pt.text = "DEVELOPED BY (TEAM 14):"
    pt.font.size = Pt(13)
    pt.font.bold = True
    pt.font.color.rgb = COLOR_GOLD
    pt.font.name = "Arial"

    members = [
        "1. Morm Sanra (ម៉ម សានរ៉ា)",
        "2. Ny Sokchan Soksdey (នី សុខច័ន្ទសួស្ដី)",
        "3. Nak Mada (ណាក់ ម៉ាដា)",
        "4. Seng Kimsant (សេង គីមសាន្ត)"
    ]
    for m in members:
        pm = tf_team.add_paragraph()
        pm.text = m
        pm.font.size = Pt(12)
        pm.font.color.rgb = COLOR_WHITE
        pm.font.name = "Arial"

    # Right Advisor Box
    tb_adv = s1.shapes.add_textbox(Inches(6.9), Inches(4.2), Inches(5.2), Inches(2.2))
    tf_adv = tb_adv.text_frame
    tf_adv.word_wrap = True
    pa = tf_adv.paragraphs[0]
    pa.text = "PROJECT ADVISOR:"
    pa.font.size = Pt(13)
    pa.font.bold = True
    pa.font.color.rgb = COLOR_GOLD
    pa.font.name = "Arial"

    pa2 = tf_adv.add_paragraph()
    pa2.text = "Prof. Lim Chonghuy (MScIT)"
    pa2.font.size = Pt(14)
    pa2.font.bold = True
    pa2.font.color.rgb = COLOR_WHITE
    pa2.font.name = "Arial"

    pa3 = tf_adv.add_paragraph()
    pa3.text = "សាស្ត្រាចារ្យ លឹម ជុងហ៊ុយ"
    pa3.font.size = Pt(13)
    pa3.font.color.rgb = COLOR_MUTED
    pa3.font.name = "Arial"

    pa4 = tf_adv.add_paragraph()
    pa4.text = "\nAcademic Year: 2026"
    pa4.font.size = Pt(12)
    pa4.font.color.rgb = COLOR_CYAN
    pa4.font.name = "Arial"

    # ── SLIDE 2: Background & Case Study ──────────────────────────────────────
    s2 = add_blank_slide()
    add_header(s2, "Project Background & Case Study Overview")

    card1 = add_card(s2, Inches(0.8), Inches(1.8), Inches(3.64), Inches(5.0), "Dara Meas Hotel")
    tb1 = s2.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(3.24), Inches(4.1))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "• Location: No. 40E, St. 2004, Phnom Penh, Cambodia.\n\n• Founded: 2019 (Established mid-scale hospitality enterprise).\n\n• Capacity: 47 fully furnished guest rooms.\n\n• Room Categories:\n  - Standard Twin\n  - Standard Double\n  - Deluxe Double\n  - Family Room\n  - Executive Suite"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE

    card2 = add_card(s2, Inches(4.84), Inches(1.8), Inches(3.64), Inches(5.0), "Operational Transition")
    tb2 = s2.shapes.add_textbox(Inches(5.04), Inches(2.5), Inches(3.24), Inches(4.1))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "• Legacy Method: Started as a family-run business relying heavily on paper logbooks and disconnected Excel spreadsheets.\n\n• Market Challenge: Rapid post-pandemic tourism growth in Phnom Penh increased competition.\n\n• Digital Imperative: Urgent need for automated management, online presence, and streamlined guest experience."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE

    card3 = add_card(s2, Inches(8.88), Inches(1.8), Inches(3.64), Inches(5.0), "Capstone Objective")
    tb3 = s2.shapes.add_textbox(Inches(9.08), Inches(2.5), Inches(3.24), Inches(4.1))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "• Modernization: Engineer a comprehensive, web-based Hotel Operations & Booking System.\n\n• Integration: Unify front-desk operations, online guest bookings, KHQR payments, and financial analytics.\n\n• Algorithmic Enhancement: Implement smart revenue management with stochastic overbooking optimization."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE

    # ── SLIDE 3: Problem Statement ─────────────────────────────────────────────
    s3 = add_blank_slide()
    add_header(s3, "Problem Statement & Legacy Operations Baseline")

    probs = [
        ("Manual Paper Bookkeeping", "Bookings recorded by hand in paper notebooks or basic Excel. High vulnerability to misplaced records, typos, and illegible entries.", COLOR_GOLD),
        ("Double-Booking & Inventory Risk", "Lack of synchronized real-time availability leads to over-booking or unfulfilled room allocations during peak tourist seasons.", COLOR_CYAN),
        ("Delayed Room State Tracking", "No immediate visibility of room states (Available / Occupied / Dirty). Housekeeping coordination happens via fragmented Telegram chats.", COLOR_GREEN),
        ("Cash-Only & Manual Auditing", "Manual payment logs cause revenue leakages, tedious nightly reconciliation, and delayed financial reporting for management.", COLOR_GOLD_LIGHT)
    ]

    for idx, (title, desc, color) in enumerate(probs):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.96)
        top = Inches(1.8 + row * 2.5)

        card = add_card(s3, left, top, Inches(5.66), Inches(2.3), border_color=color)
        tb = s3.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.26), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{idx+1}. {title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_WHITE

    # ── SLIDE 4: Objectives & Scope ────────────────────────────────────────────
    s4 = add_blank_slide()
    add_header(s4, "Project Objectives & Boundary Scope")

    c1 = add_card(s4, Inches(0.8), Inches(1.8), Inches(5.66), Inches(5.0), "Core Objectives", border_color=COLOR_GREEN)
    tb1 = s4.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.26), Inches(4.1))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "1. Operational Automation: Replace paper ledgers with a centralized digital management system.\n\n2. Real-Time Room Tracking: Provide instant status updates across Front Desk, Housekeeping, and Admin.\n\n3. Integrated Cashless Payments: Enable Bakong KHQR & ABA Payway QR code payment processing.\n\n4. Smart Revenue Optimization: Implement an automated overbooking multiplier heuristic to maximize yield."
    p1.font.size = Pt(13)
    p1.font.color.rgb = COLOR_WHITE

    c2 = add_card(s4, Inches(6.86), Inches(1.8), Inches(5.66), Inches(5.0), "Scope Boundaries Matrix", border_color=COLOR_CYAN)
    tb2 = s4.shapes.add_textbox(Inches(7.06), Inches(2.5), Inches(5.26), Inches(4.1))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "IN-SCOPE:\n• Multi-role portal (Admin, Manager, Receptionist, Guest).\n• Guest online booking & direct search filter.\n• Walk-in proxy booking & stay extension.\n• KHQR / ABA Payway payment confirmation.\n• Overbooking optimization command & analytics.\n\nOUT-OF-SCOPE:\n• Global OTA API sync (Booking.com / Airbnb).\n• Full enterprise ERP/accounting software.\n• IoT smart door lock hardware integration."
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_WHITE

    # ── SLIDE 5: Development Methodology ──────────────────────────────────────
    s5 = add_blank_slide()
    add_header(s5, "Software Development Life Cycle: Agile Framework")

    steps = [
        ("1. Exploration", "Req. gathering, staff interviews, current workflow mapping.", COLOR_GOLD),
        ("2. Planning", "Architecture design, User Stories, Planning Game, DFDs.", COLOR_CYAN),
        ("3. Iterations (Sprints)", "2-3 week sprints, feature development, unit/integration testing.", COLOR_GREEN),
        ("4. Productionizing", "Local server setup, XAMPP/MySQL config, system deployment.", COLOR_GOLD_LIGHT),
        ("5. Maintenance", "System auditing, bug fixes, data backup & user feedback.", COLOR_GOLD)
    ]

    for idx, (title, desc, color) in enumerate(steps):
        left = Inches(0.8 + idx * 2.38)
        top = Inches(1.8)
        width = Inches(2.2)
        height = Inches(5.0)

        card = add_card(s5, left, top, width, height, border_color=color)
        tb = s5.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        pd = tf.add_paragraph()
        pd.text = f"\n{desc}"
        pd.font.size = Pt(12)
        pd.font.color.rgb = COLOR_WHITE

    # ── SLIDE 6: System Architecture & RBAC ────────────────────────────────────
    s6 = add_blank_slide()
    add_header(s6, "Role-Based Access Control (RBAC) Architecture")

    roles = [
        ("Administrator", "Full system control, staff management, room pricing config, backup/restore, overbooking parameters.", COLOR_GOLD),
        ("Manager", "Business logic overview, revenue reporting, occupancy analytics, cancellation rates analysis.", COLOR_CYAN),
        ("Receptionist", "Front-desk operations, walk-in registration, check-in/out processing, stay extensions, KHQR payment verification.", COLOR_GREEN),
        ("Guest / Public User", "Online room browsing, date availability filter, self-service booking, QR payment submission.", COLOR_GOLD_LIGHT)
    ]

    for idx, (title, desc, color) in enumerate(roles):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.96)
        top = Inches(1.8 + row * 2.5)

        card = add_card(s6, left, top, Inches(5.66), Inches(2.3), border_color=color)
        tb = s6.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.26), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Role: {title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        p_desc = tf.add_paragraph()
        p_desc.text = f"\nKey Capabilities:\n{desc}"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_WHITE

    # ── SLIDE 7: DFD & System Diagrams ───────────────────────────────────────
    s7 = add_blank_slide()
    add_header(s7, "System Analysis: Data Flow Diagram (DFD)")

    c1 = add_card(s7, Inches(0.8), Inches(1.8), Inches(5.66), Inches(5.0), "Context Diagram & Entities", border_color=COLOR_GOLD)
    tb1 = s7.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.26), Inches(4.1))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "External Entities:\n• Guest: Booking requests, payment details, feedback.\n• Receptionist: Check-in/out triggers, proxy bookings.\n• Admin: Rate changes, user management.\n• Payment Gateway: KHQR payment status callback.\n\nKey Core Processes (DFD Level 0):\n1.0 Web Access & Room Search\n2.0 Payment Processing\n3.0 Check-In / Check-Out Operations\n4.0 Staff & Access Administration\n5.0 Stay Extension\n6.0 Room Services Management\n7.0 System Administration"
    p1.font.size = Pt(12)
    p1.font.color.rgb = COLOR_WHITE

    c2 = add_card(s7, Inches(6.86), Inches(1.8), Inches(5.66), Inches(5.0), "Data Stores & State Machine", border_color=COLOR_CYAN)
    tb2 = s7.shapes.add_textbox(Inches(7.06), Inches(2.5), Inches(5.26), Inches(4.1))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Primary Data Stores:\n• D1: Guests (Profile, contact credentials)\n• D2: Rooms (Room number, status, capacity)\n• D3: Bookings (Dates, statuses, multipliers)\n• D4: Transactions (Payments, receipts, KHQR ref)\n• D5: Staffs (User credentials, RBAC roles)\n• D6: Admins (System configurations)\n\nRoom State Transitions:\nAvailable ──(Check-In)──> Occupied ──(Check-Out)──> Dirty ──(Housekeeping)──> Available"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_WHITE

    # ── SLIDE 8: Database Architecture ────────────────────────────────────────
    s8 = add_blank_slide()
    add_header(s8, "Database Architecture & Entity-Relationship Schema")

    card = add_card(s8, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), "Relational Database Structure (MySQL)", border_color=COLOR_GOLD)

    tb = s8.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Core Tables & Key Relationships:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD

    schema_info = (
        "• room_types: (id, name, price_per_night, overbooking_multiplier, capacity)\n"
        "• rooms: (id, room_number, room_type_id [FK], current_status: available/occupied/dirty)\n"
        "• guests: (id, full_name, email, phone, gender, nationality)\n"
        "• bookings: (id, guest_id [FK], room_id [FK], check_in_date, check_out_date, booking_status: pending/booked/checked_in/checked_out/no_show/relocated)\n"
        "• transactions: (id, booking_id [FK], amount_paid, payment_method: cash/khqr, payment_status: paid/partial/refunded)\n"
        "• staff: (id, full_name, username, role: admin/manager/receptionist, password_hash)\n"
        "• room_services: (id, booking_id [FK], request_type, status, notes)"
    )

    p_body = tf.add_paragraph()
    p_body.text = f"\n{schema_info}"
    p_body.font.size = Pt(13)
    p_body.font.color.rgb = COLOR_WHITE

    # ── SLIDE 9: Core Features & Payment Gateway ──────────────────────────────
    s9 = add_blank_slide()
    add_header(s9, "Core Features & Cashless Payment Integration")

    features = [
        ("Online Guest Self-Booking", "Guests can browse available rooms, filter by dates and room types, select options, and submit booking requests 24/7.", COLOR_GOLD),
        ("KHQR & ABA Payway Integration", "Generates dynamic KHQR payment codes via Bakong API & ABA Payway. Instant verification and electronic receipts.", COLOR_CYAN),
        ("Front Desk Desk Portal", "Receptionists easily handle walk-in bookings, check-in, check-out, balance collection, and stay extensions in seconds.", COLOR_GREEN),
        ("Analytical Management Dashboard", "Admins view real-time occupancy rates, monthly revenue graphs, booking logs, and automatic data backups.", COLOR_GOLD_LIGHT)
    ]

    for idx, (title, desc, color) in enumerate(features):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.96)
        top = Inches(1.8 + row * 2.5)

        card = add_card(s9, left, top, Inches(5.66), Inches(2.3), border_color=color)
        tb = s9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.26), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_WHITE

    # ── SLIDE 10: Algorithmic Overbooking Optimization ────────────────────────
    s10 = add_blank_slide()
    add_header(s10, "Algorithmic Innovation: Overbooking Optimization")

    c1 = add_card(s10, Inches(0.8), Inches(1.8), Inches(5.66), Inches(5.0), "Stochastic Heuristic Model", border_color=COLOR_GOLD)
    tb1 = s10.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.26), Inches(4.1))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Command: app:optimize-overbooking\nSchedule: Nightly at 02:05 AM (after night audit)\nTheoretical Basis: Talluri & van Ryzin Revenue Management\n\nSelf-Tuning Multiplier Logic:\n• Penalty Step (-0.05):\n  Applied if relocations > 0 (overbooking was too aggressive; dial back risk).\n• Reward Step (+0.01):\n  Applied if no-shows > 0 and relocations = 0 (played too safe; rooms went wasted).\n• Neutral Step (0.00):\n  No relocations and no no-shows.\n\nClamped Range: 1.00 ≤ overbooking_multiplier ≤ 1.50"
    p1.font.size = Pt(12)
    p1.font.color.rgb = COLOR_WHITE

    c2 = add_card(s10, Inches(6.86), Inches(1.8), Inches(5.66), Inches(5.0), "Code Implementation Highlight", border_color=COLOR_CYAN)
    tb2 = s10.shapes.add_textbox(Inches(7.06), Inches(2.5), Inches(5.26), Inches(4.1))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Code Snippet (OptimizeOverbooking.php):\n\n$relocations = $yesterdayBookings\n  ->where('status', STATUS_RELOCATED)->count();\n$noShows = $yesterdayBookings\n  ->where('status', STATUS_NO_SHOW)->count();\n\nif ($relocations > 0) {\n    $adjusted -= self::PENALTY_STEP; // -0.05\n} elseif ($noShows > 0) {\n    $adjusted += self::REWARD_STEP;  // +0.01\n}\n\n$roomType->overbooking_multiplier =\n    min(1.50, max(1.00, $adjusted));"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_WHITE

    # ── SLIDE 11: Technology Stack ────────────────────────────────────────────
    s11 = add_blank_slide()
    add_header(s11, "Technology Stack & Technical Infrastructure")

    techs = [
        ("Backend Framework", "• PHP 8.x / 11\n• Laravel Framework (MVC)\n• Composer Package Manager\n• Artisan CLI Commands", COLOR_GOLD),
        ("Frontend & UI", "• React.js / Blade Views\n• Vanilla CSS & Tailwind CSS\n• Bootstrap Icons\n• Responsive Design", COLOR_CYAN),
        ("Database & Web Server", "• MySQL RDBMS\n• Apache Web Server\n• XAMPP Control Panel\n• Windows 10/11 Deployment", COLOR_GREEN),
        ("DevOps & Utilities", "• Git / GitHub Version Control\n• VS Code IDE\n• Postman API Tester\n• ReportLab / PDF Utilities", COLOR_GOLD_LIGHT)
    ]

    for idx, (title, desc, color) in enumerate(techs):
        left = Inches(0.8 + idx * 2.98)
        top = Inches(1.8)
        width = Inches(2.78)
        height = Inches(5.0)

        card = add_card(s11, left, top, width, height, border_color=color)
        tb = s11.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        pd = tf.add_paragraph()
        pd.text = f"\n{desc}"
        pd.font.size = Pt(12)
        pd.font.color.rgb = COLOR_WHITE

    # ── SLIDE 12: UI Demonstration ───────────────────────────────────────────
    s12 = add_blank_slide()
    add_header(s12, "System Interface & Portal Overview")

    portals = [
        ("Guest Portal", "Modern homepage, interactive room catalog, date availability picker, online reservation form, instant KHQR code modal.", COLOR_GOLD),
        ("Reception Portal", "Front-desk operations board, quick walk-in guest check-in, real-time room availability grid, stay extension popup, check-out receipt printer.", COLOR_CYAN),
        ("Admin Portal", "Financial performance dashboard, monthly revenue graphs, room type price adjustment, staff account creation, system backup controls.", COLOR_GREEN)
    ]

    for idx, (title, desc, color) in enumerate(portals):
        left = Inches(0.8 + idx * 3.98)
        top = Inches(1.8)
        width = Inches(3.78)
        height = Inches(5.0)

        card = add_card(s12, left, top, width, height, border_color=color)
        tb = s12.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        pd = tf.add_paragraph()
        pd.text = f"\n{desc}"
        pd.font.size = Pt(13)
        pd.font.color.rgb = COLOR_WHITE

    # ── SLIDE 13: Results & Impact ────────────────────────────────────────────
    s13 = add_blank_slide()
    add_header(s13, "Operational Results & Business Impact Evaluation")

    metrics = [
        ("+60%", "Staff Workload Efficiency", "Automated booking and room tracking drastically reduced front-desk paperwork.", COLOR_GREEN),
        ("-90%", "Reduction in Errors", "Eliminated paper ledger mistakes and prevented double-booking risks.", COLOR_GOLD),
        ("-50%", "Guest Check-In Time", "Streamlined KHQR payment & instant check-in reduced lobby waiting time.", COLOR_CYAN),
        ("100%", "Data Auditability", "Full digital audit trail for all transactions, bookings, and room status updates.", COLOR_GOLD_LIGHT)
    ]

    for idx, (stat, title, desc, color) in enumerate(metrics):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.96)
        top = Inches(1.8 + row * 2.5)

        card = add_card(s13, left, top, Inches(5.66), Inches(2.3), border_color=color)
        tb = s13.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.26), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p_stat = tf.paragraphs[0]
        p_stat.text = f"{stat}  —  {title}"
        p_stat.font.size = Pt(18)
        p_stat.font.bold = True
        p_stat.font.color.rgb = color

        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_WHITE

    # ── SLIDE 14: Future Enhancements & Conclusion ────────────────────────────
    s14 = add_blank_slide()
    add_header(s14, "Future Roadmap & Defense Conclusion")

    c1 = add_card(s14, Inches(0.8), Inches(1.8), Inches(5.66), Inches(5.0), "Future System Roadmap", border_color=COLOR_GOLD)
    tb1 = s14.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5.26), Inches(4.1))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "1. Global Channel Manager:\nIntegrate Booking.com and Agoda API for automated two-way inventory syncing.\n\n2. Smart IoT Door Locks:\nImplement Bluetooth / QR Code smart lock access for keyless guest check-in.\n\n3. Restaurant POS Module:\nConnect hotel restaurant & room service billing directly to guest room folios."
    p1.font.size = Pt(13)
    p1.font.color.rgb = COLOR_WHITE

    c2 = add_card(s14, Inches(6.86), Inches(1.8), Inches(5.66), Inches(5.0), "Conclusion & Q&A Session", border_color=COLOR_CYAN)
    tb2 = s14.shapes.add_textbox(Inches(7.06), Inches(2.5), Inches(5.26), Inches(4.1))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "THANK YOU!\n\nThe Dara Meas Hotel Booking Operation System successfully fulfills all Bachelor Capstone Requirements for Norton University.\n\nWe would like to thank our Advisor, Prof. Lim Chonghuy, and the Honorable Examination Committee.\n\nWe are now ready for your questions & feedback."
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), "Dara_Meas_Hotel_Capstone_Defense.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_deck()
